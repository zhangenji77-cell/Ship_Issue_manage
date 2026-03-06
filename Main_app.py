import time
import re
import io
import os
import zipfile
import tempfile
import subprocess
from datetime import datetime, timedelta
import pandas as pd
import sqlalchemy
from sqlalchemy import text
import streamlit as st
import openpyxl
from openpyxl.styles import Alignment, Font, Border, Side
from pptx import Presentation
from pptx.util import Inches, Pt as Ppt_Pt
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Pt, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_LINE_SPACING

# --- 1. Basic Configuration & CSS ---
st.set_page_config(page_title="TSM Summary of Weekly Ship Reports", layout="wide")

st.markdown("""
    <style>
    .stButton>button { width: 100%; border-radius: 5px; height: 3em; }
    .stDownloadButton>button { width: 100%; border-radius: 5px; background-color: #004a99; color: white; }
    div.stButton > button[key^="import_"] {
        background-color: #f8f9fa !important;
        color: #004a99 !important;
        border: 1px solid #004a99 !important;
    }
    </style>
    """, unsafe_allow_html=True)

# Initialize Session State
if 'logged_in' not in st.session_state: st.session_state.logged_in = False
if 'username' not in st.session_state: st.session_state.username = None
if 'role' not in st.session_state: st.session_state.role = None
if 'ship_index' not in st.session_state: st.session_state.ship_index = 0
if 'drafts' not in st.session_state: st.session_state.drafts = {}
if 'editing_id' not in st.session_state: st.session_state.editing_id = None
if 'confirm_del_id' not in st.session_state: st.session_state.confirm_del_id = None

@st.cache_resource
def get_engine():
    return sqlalchemy.create_engine(st.secrets["postgres_url"])


# --- 2. Report Generation Tools ---

def generate_custom_excel(df):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Ship Report"

    font_yahei = Font(name='微软雅黑', size=10)
    font_yahei_bold = Font(name='微软雅黑', size=10, bold=True)
    thin_side = Side(style='thin', color='000000')
    black_border = Border(top=thin_side, left=thin_side, right=thin_side, bottom=thin_side)

    ws.merge_cells('A1:C1')
    ws['A1'] = f"Report Date: {datetime.now().strftime('%Y-%m-%d')}"
    ws['A1'].font = Font(name='微软雅黑', size=12, bold=True)
    ws['A1'].alignment = Alignment(horizontal='center', vertical='center')

    headers = ['Manager Name', 'Vessel Name', 'Issue']
    for col_num, header in enumerate(headers, 1):
        cell = ws.cell(row=2, column=col_num, value=header)
        cell.font = font_yahei_bold
        cell.alignment = Alignment(horizontal='center', vertical='center')
        cell.border = black_border

    def clean_and_reformat_issue(series):
        all_lines = []
        for content in series:
            if content:
                lines = str(content).split('\n')
                for line in lines:
                    clean_line = re.sub(r'^\d+[\.、\s]*', '', line.strip())
                    if clean_line:
                        all_lines.append(clean_line)
        if not all_lines: return ""
        return "\n".join([f"{i + 1}. {text}" for i, text in enumerate(all_lines)])

    df_grouped = df.groupby(['manager_name', 'ship_name'])['this_week_issue'].apply(
        clean_and_reformat_issue).reset_index()
    df_grouped = df_grouped.sort_values(by='manager_name')

    current_row = 3
    for manager, group in df_grouped.groupby('manager_name', sort=False):
        start_merge_row = current_row
        for _, row_data in group.iterrows():
            for col in [1, 2]:
                cell = ws.cell(row=current_row, column=col,
                               value=row_data['manager_name'] if col == 1 else row_data['ship_name'])
                cell.font = font_yahei
                cell.border = black_border
                cell.alignment = Alignment(horizontal='center', vertical='center')

            cell_c = ws.cell(row=current_row, column=3, value=row_data['this_week_issue'])
            cell_c.font = font_yahei
            cell_c.border = black_border
            cell_c.alignment = Alignment(wrap_text=True, horizontal='left', vertical='center')
            current_row += 1

        if len(group) > 1:
            ws.merge_cells(start_row=start_merge_row, start_column=1, end_row=current_row - 1, end_column=1)
            for r in range(start_merge_row, current_row):
                ws.cell(row=r, column=1).border = black_border

    ws.column_dimensions['A'].width = 20
    ws.column_dimensions['B'].width = 25
    ws.column_dimensions['C'].width = 70

    output = io.BytesIO()
    wb.save(output)
    return output.getvalue()


def create_ppt_report(df, start_date, end_date):
    prs = Presentation()

    slide_layout_title = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout_title)

    try:
        slide.shapes.add_picture("TSM_Logo.png", left=Inches(4.25), top=Inches(1.2), width=Inches(1.8))
    except:
        pass

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.top = Inches(3.5)
    title.text = "TSM Summary of Weekly Ship Reports"

    current_date = datetime.now().strftime('%Y-%m-%d')
    subtitle.text = f"Creation Date: {current_date}"
    subtitle.top = Inches(4.5)

    def clean_and_reformat_ppt(series):
        all_lines = []
        for content in series:
            if content:
                lines = str(content).split('\n')
                for line in lines:
                    clean_line = re.sub(r'^\d+[\.、\s]*', '', line.strip())
                    if clean_line: all_lines.append(clean_line)
        if not all_lines: return ""
        return "\n".join([f"{i + 1}. {text}" for i, text in enumerate(all_lines)])

    df_ppt = df.groupby(['manager_name', 'ship_name'], sort=False)['this_week_issue'].apply(
        clean_and_reformat_ppt).reset_index()
    df_ppt = df_ppt.sort_values(by=['manager_name', 'ship_name'])

    for _, row in df_ppt.iterrows():
        manager = row['manager_name']
        ship = row['ship_name']
        issue_content = row['this_week_issue']

        slide_layout_content = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout_content)

        slide.shapes.title.text = f" {ship} ({manager})"
        tf = slide.placeholders[1].text_frame
        tf.word_wrap = True

        if issue_content:
            for line in issue_content.split('\n'):
                p = tf.add_paragraph()
                p.text = line
                p.level = 0
                p.font.size = Pt(24)
                p.font.name = '微软雅黑'

    slide_layout_blank = prs.slide_layouts[6]
    end_slide = prs.slides.add_slide(slide_layout_blank)
    tx_box = end_slide.shapes.add_textbox(Inches(3), Inches(3.5), Inches(4), Inches(2))
    tf_end = tx_box.text_frame
    tf_end.text = "Thank you for watching."
    p_end = tf_end.paragraphs[0]
    p_end.alignment = PP_ALIGN.CENTER
    p_end.font.size = Pt(44)
    p_end.font.bold = True
    p_end.font.name = '微软雅黑'

    ppt_out = io.BytesIO()
    prs.save(ppt_out)
    ppt_out.seek(0)
    return ppt_out


# ---------------------------------------------------------
# Paylist Generator Logic (工资单生成逻辑)
# ---------------------------------------------------------
def normalize_key(key):
    if pd.isna(key): return ""
    return re.sub(r'\s+', '', str(key)).lower()


def clean_filename(name):
    return re.sub(r'[\\/*?:"<>|]', "", str(name)).strip()


def format_currency(val):
    if pd.isna(val) or val == "": return ""
    try:
        s_val = str(val).replace(',', '').strip()
        if not s_val: return ""
        return "{:,.2f}".format(float(s_val))
    except (ValueError, TypeError):
        return str(val)


def format_date_custom(val):
    if pd.isna(val) or val == "": return ""
    try:
        if hasattr(val, 'strftime'): return val.strftime('%d/%m/%Y')
        s_val = str(val).split()[0].strip()
        if '-' in s_val:
            parts = s_val.split('-')
            if len(parts) == 3 and len(parts[0]) == 4:
                return f"{parts[2]}/{parts[1]}/{parts[0]}"
        return s_val
    except:
        return str(val)


def set_cell_text(cell, text):
    if text is None: text = ""
    text = str(text)
    if text.endswith(".0"): text = text[:-2]
    cell.text = ""
    p = cell.paragraphs[0]
    run = p.add_run(text)
    run.font.size = Pt(9)
    run.font.name = 'Arial Narrow'
    run.font.bold = True
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after = Pt(0)
    p.paragraph_format.line_spacing = 1.5


def shrink_empty_lines(doc):
    for p in doc.paragraphs:
        if not p.text.strip():
            p_fmt = p.paragraph_format
            p_fmt.space_before = Pt(0)
            p_fmt.space_after = Pt(0)
            p_fmt.line_spacing = 1.0
            if p.runs:
                for r in p.runs: r.font.size = Pt(1)
            else:
                p.add_run(" ").font.size = Pt(1)


def insert_spacer_before_payslip(doc):
    for p in doc.paragraphs:
        if "PAY SLIP" in p.text:
            spacer = p.insert_paragraph_before(" ")
            spacer.paragraph_format.space_after = Pt(0)
            spacer.paragraph_format.line_spacing = 1.0
            # 💡 将字体大小从 Pt(12) 增大到 Pt(36)，利用这个隐藏的空行把标题往下挤
            if spacer.runs:
                spacer.runs[0].font.size = Pt(12)
            else:
                spacer.add_run(" ").font.size = Pt(12)
            break


def generate_paylist_zip(uploaded_excel):
    """读取上传的 Excel，生成包含所有 Word 工资单的 ZIP 压缩包"""
    # 1. 从内存中读取上传的文件
    df_raw = pd.read_excel(uploaded_excel, sheet_name='SUM-SAL', header=None)

    employees = []
    current_vessel = "Unknown Vessel"
    i = 0

    # 2. 提取数据 (保留了您原本的清洗逻辑)
    while i < len(df_raw):
        row = df_raw.iloc[i].tolist()
        first_cell = str(row[0]).strip() if pd.notna(row[0]) else ""

        if i + 1 < len(df_raw):
            next_row_first = str(df_raw.iloc[i + 1][0]).strip()
            if next_row_first == 'S/N':
                if "Vessel Name:" in first_cell:
                    current_vessel = row[1]
                elif first_cell and first_cell.lower() != 'nan':
                    current_vessel = first_cell
                raw_headers = df_raw.iloc[i + 1].tolist()
                headers_map = {normalize_key(h): idx for idx, h in enumerate(raw_headers) if pd.notna(h)}
                i += 2
                continue

        if 's/n' not in locals().get('headers_map', {}) and 'name' not in locals().get('headers_map', {}):
            i += 1;
            continue

        if first_cell.isdigit():
            emp = {'Vessel Name': current_vessel}

            def get_val(col_keywords):
                for key in headers_map:
                    if normalize_key(col_keywords) in key:
                        val = row[headers_map[key]]
                        return val if pd.notna(val) else ""
                return ""

            emp['Name'] = get_val('Name')
            emp['Rank'] = get_val('Rank')
            emp['From'] = format_date_custom(get_val('From(Date)'))
            emp['To'] = format_date_custom(get_val('To(Date)'))
            emp['Day on Board'] = get_val('Day on Board')
            emp['Basic Salary'] = format_currency(get_val('Basic Salary'))
            emp['Fixed OT'] = format_currency(get_val('Fixed OT'))
            emp['Leave Pay'] = format_currency(get_val('Leave Pay'))
            emp['Allowance'] = format_currency(get_val('Allowance'))
            emp['Net Salary'] = format_currency(get_val('Net Salary'))
            emp['Reimbursement'] = format_currency(get_val('Reimbursement'))
            emp['Subtotal'] = format_currency(get_val('Subtotal'))
            emp['Deduction'] = format_currency(get_val('Deduction'))
            emp['Release'] = format_currency(get_val('Release'))
            emp['Retaining'] = format_currency(get_val('Retaining'))

            rem_foreign = get_val('Remittance - Foreign')
            rem_sg = get_val('Remittance - Singapore')
            use_foreign = False
            try:
                if rem_foreign and float(str(rem_foreign).replace(',', '')) > 0: use_foreign = True
            except:
                pass

            emp['Remittance'] = format_currency(rem_foreign if use_foreign else rem_sg)
            emp['Remarks'] = get_val('Remarks')
            employees.append(emp)
        i += 1

    # 3. 生成 Word 文档并压缩进 ZIP (纯内存操作，速度极快)
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
        for emp in employees:
            # ⚠️ 必须确保 payslip模版.docx 文件存在于服务器目录下
            doc = Document('payslip模版.docx')
            insert_spacer_before_payslip(doc)

            section = doc.sections[0]
            section.top_margin, section.bottom_margin = Cm(1.0), Cm(0.5)
            section.left_margin, section.right_margin = Cm(1.0), Cm(1.0)
            tables = doc.tables

            def fill_simple(table, label, value):
                for row in table.rows:
                    for c, cell in enumerate(row.cells):
                        if label in cell.text and c + 1 < len(row.cells):
                            set_cell_text(row.cells[c + 1], value)
                            return

            fill_simple(tables[0], "Employee's Name", emp['Name'])
            fill_simple(tables[0], "Vessel Name", emp['Vessel Name'])
            fill_simple(tables[1], "Rank", emp['Rank'])
            fill_simple(tables[1], "FROM", emp['From'])
            fill_simple(tables[1], "TO", emp['To'])
            fill_simple(tables[1], "Day on Board", emp['Day on Board'])

            t2 = tables[2]
            header_row_idx, col_earn, col_deduct = -1, -1, -1
            for r_idx in range(min(5, len(t2.rows))):
                amount_indices = [c_idx for c_idx, cell in enumerate(t2.rows[r_idx].cells) if 'Amount' in cell.text]
                if len(amount_indices) >= 2:
                    header_row_idx, col_earn, col_deduct = r_idx, amount_indices[0], amount_indices[-1]
                    break

            if col_earn != -1 and col_deduct != -1:
                def fill_left(label, val):
                    for r in range(header_row_idx + 1, len(t2.rows)):
                        if normalize_key(label) in normalize_key(
                                "".join([c.text for c in t2.rows[r].cells[:col_earn]])):
                            set_cell_text(t2.rows[r].cells[col_earn], val)
                            break

                def fill_right(label, val):
                    for r in range(header_row_idx + 1, len(t2.rows)):
                        if normalize_key(label) in normalize_key("".join([c.text for c in t2.rows[r].cells])):
                            set_cell_text(t2.rows[r].cells[col_deduct], val)
                            break

                fill_left('Basic Salary', emp['Basic Salary'])
                fill_left('Fixed OT', emp['Fixed OT'])
                fill_left('Leave Pay', emp['Leave Pay'])
                fill_left('Allowance', emp['Allowance'])
                fill_left('Total Earnings', emp['Net Salary'])
                fill_left('Reimbursement', emp['Reimbursement'])
                fill_left('Net Amount', emp['Subtotal'])
                fill_right('Total Deductions', emp['Deduction'])
                fill_right('Release', emp['Release'])
                fill_right('Retaining', emp['Retaining'])
                fill_right('Remittance', emp['Remittance'])

            remarks_content = str(emp['Remarks']).strip()
            if remarks_content and remarks_content.lower() != 'nan' and remarks_content != '0':
                for p in doc.paragraphs:
                    if "Remarks:" in p.text:
                        run = p.add_run(" " + remarks_content)
                        run.font.size, run.font.name, run.font.bold = Pt(9), 'Arial Narrow', False
                        p.paragraph_format.line_spacing = 1.0
                        break

            shrink_empty_lines(doc)

            # 将单个文档保存到内存
            doc_buffer = io.BytesIO()
            doc.save(doc_buffer)
            doc_buffer.seek(0)

            # 写入 ZIP 文件中（自动按船名建立文件夹）
            safe_vessel = clean_filename(emp['Vessel Name']) or "Uncategorized"
            safe_emp = clean_filename(emp['Name'])
            zip_file.writestr(f"{safe_vessel}/{safe_emp}.docx", doc_buffer.getvalue())

    zip_buffer.seek(0)
    return zip_buffer


# =========================================================
# 新增功能：进阶版 Paylist 生成逻辑 (动态计算 + Word + PDF 双版本)
# =========================================================
def generate_advanced_paylist_zip(uploaded_excel):
    """读取上传的 Excel，动态计算薪资，并在安全屋中生成 Word 和 PDF 双版本 ZIP 压缩包"""
    # 每次调用时将指针重置到开头
    uploaded_excel.seek(0)

    # 1. 智能查找目标 Sheet
    xl = pd.ExcelFile(uploaded_excel)
    target_sheet = None
    for sheet in xl.sheet_names:
        # 忽略大小写、去掉空格进行匹配
        if 'SUM-SAL' in sheet.upper().replace(' ', ''):
            target_sheet = sheet
            break

    if not target_sheet:
        raise ValueError(f"未找到包含 'SUM-SAL' 的工作表！当前文件包含的表有: {', '.join(xl.sheet_names)}")

    # 2. 从内存中读取锁定的 Sheet
    df_raw = pd.read_excel(xl, sheet_name=target_sheet, header=None)

    employees = []
    current_vessel = "Unknown Vessel"
    headers_map = {}
    i = 0
    # ... (下方保留你原有的 while i < len(df_raw): 代码不变) ...

    while i < len(df_raw):
        row_vals = [str(x).strip() for x in df_raw.iloc[i].tolist()]
        if 'S/N' in row_vals:
            if i > 0:
                prev_row = [str(x).strip() for x in df_raw.iloc[i - 1].tolist() if
                            pd.notna(x) and str(x).strip() not in ['', 'nan']]
                if prev_row:
                    v_name = prev_row[0]
                    current_vessel = v_name.split(":", 1)[1].strip() if "Vessel Name:" in v_name else v_name
            headers_map = {normalize_key(h): idx for idx, h in enumerate(row_vals) if h not in ['nan', '']}
            i += 1;
            continue

        first_cell = str(df_raw.iloc[i][0]).strip()
        if first_cell.isdigit() and headers_map:
            row_data = df_raw.iloc[i].tolist()

            def get_val(col_keywords):
                norm_key = normalize_key(col_keywords)
                for key, idx in headers_map.items():
                    if norm_key in key:
                        return row_data[idx] if pd.notna(row_data[idx]) else ""
                return ""

            try:
                m_str = str(get_val('MonthlySalary')).replace(',', '').strip()
                m_val = float(m_str) if m_str else 0.0
            except:
                m_val = 0.0

            basic_val = m_val * 0.58
            ot_val = m_val * 0.37
            fixed_ot = int(ot_val) + 1 if (round(ot_val, 4) - int(ot_val)) >= 0.5 else int(ot_val)
            lp_val = m_val * 0.05
            leave_pay = int(lp_val) + 1 if (round(lp_val, 4) - int(lp_val)) > 0.5 else int(lp_val)

            emp = {
                'Vessel Name': current_vessel, 'Name': get_val('Name'), 'Rank': get_val('Rank'),
                'From': format_date_custom(get_val('FromDate') or get_val('From')),
                'To': format_date_custom(get_val('ToDate') or get_val('To')),
                'Day on Board': str(get_val('DayonBoard')),
                'Basic Salary': format_currency(basic_val), 'Fixed OT': format_currency(fixed_ot),
                'Leave Pay': format_currency(leave_pay), 'Total Earnings': format_currency(m_val),
                'Reimbursement': format_currency(get_val('Reimbursement')),
                'Net Amount': format_currency(get_val('SubTotal')),
                'Total Deductions': format_currency(get_val('Deduction')),
                'Release': format_currency(get_val('ReleaseofSalary')),
                'Retaining': format_currency(get_val('Retaining')),
                'Remittance': format_currency(get_val('RemittanceForeignBank') or get_val('Remittance')),
                'Remarks': get_val('Remarks')
            }
            employees.append(emp)
        i += 1

    zip_buffer = io.BytesIO()

    # 开启安全屋，利用 LibreOffice 生成 PDF
    with tempfile.TemporaryDirectory() as temp_dir:
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
            for emp in employees:
                # ⚠️ 确保服务器里上传了这个新模版文件
                doc = Document('Out_port paylist 模版.docx')
                insert_spacer_before_payslip(doc)

                section = doc.sections[0]
                section.top_margin, section.bottom_margin = Cm(1.0), Cm(0.5)
                section.left_margin, section.right_margin = Cm(1.0), Cm(1.0)

                def fill_simple(table, label, value):
                    for row in table.rows:
                        for c, cell in enumerate(row.cells):
                            txt = cell.text.strip()
                            if label in ["TO", "FROM", "Rank"] and txt not in [label, f"{label}:",
                                                                               f"{label} :"]: continue
                            if label in cell.text and c + 1 < len(row.cells):
                                set_cell_text(row.cells[c + 1], value)
                                return

                for table in doc.tables[:2]:
                    fill_simple(table, "Employee's Name", emp['Name'])
                    fill_simple(table, "Vessel Name", emp['Vessel Name'])
                    fill_simple(table, "Rank", emp['Rank'])
                    fill_simple(table, "FROM", emp['From'])
                    fill_simple(table, "TO", emp['To'])
                    fill_simple(table, "Day on Board", emp['Day on Board'])

                if len(doc.tables) >= 3:
                    t_fin = doc.tables[2]
                    col_earn, col_deduct = 4, 9
                    for r in range(min(5, len(t_fin.rows))):
                        amts = [idx for idx, c in enumerate(t_fin.rows[r].cells) if 'Amount' in c.text]
                        if len(amts) >= 2:
                            col_earn, col_deduct = amts[0], amts[-1]
                            break

                    for row in t_fin.rows:
                        label = row.cells[0].text.strip()
                        if col_earn < len(row.cells):
                            if "Basic Salary" in label:
                                set_cell_text(row.cells[col_earn], emp['Basic Salary'])
                            elif "Fixed OT" in label:
                                set_cell_text(row.cells[col_earn], emp['Fixed OT'])
                            elif "Leave Pay" in label:
                                set_cell_text(row.cells[col_earn], emp['Leave Pay'])
                            elif "Total Earnings" in label:
                                set_cell_text(row.cells[col_earn], emp['Total Earnings'])
                            elif "Reimbursement" in label:
                                set_cell_text(row.cells[col_earn], emp['Reimbursement'])
                            elif "Net Amount" in label:
                                set_cell_text(row.cells[col_earn], emp['Net Amount'])

                        for idx, cell in enumerate(row.cells):
                            c_txt = cell.text.strip()
                            if col_deduct < len(row.cells):
                                if "Total Deductions" in c_txt:
                                    set_cell_text(row.cells[col_deduct], emp['Total Deductions'])
                                elif "Release" in c_txt:
                                    set_cell_text(row.cells[col_deduct], emp['Release'])
                                elif "Retaining" in c_txt:
                                    set_cell_text(row.cells[col_deduct], emp['Retaining'])
                                elif "Remittance - Bank" in c_txt:
                                    set_cell_text(row.cells[col_deduct], emp['Remittance'])

                rem = str(emp['Remarks']).strip()
                if rem and rem.lower() != 'nan' and rem != '0':
                    for p in doc.paragraphs:
                        if "Remarks:" in p.text:
                            run = p.add_run(" " + rem)
                            run.font.size, run.font.name, run.font.bold = Pt(9), 'Arial Narrow', False
                            p.paragraph_format.line_spacing = 1.0
                            break

                shrink_empty_lines(doc)

                safe_vessel = clean_filename(emp['Vessel Name']) or "Uncategorized"
                safe_emp = clean_filename(emp['Name'])

                # 保存为 Word
                # 1. 首先，将排版正常的【纯 Word 版】原封不动地保存下来
                temp_docx_path = os.path.join(temp_dir, f"{safe_emp}.docx")
                doc.save(temp_docx_path)

                # 2. ⚡️ 核心魔法：专门为 PDF 版本单独增加向下的间距 ⚡️
                # 扫描文档，找到 "PAY SLIP" 这几个字，强行给它上面加空白
                for p in doc.paragraphs:
                    if "PAY SLIP" in p.text:
                        # 这里的 Pt(30) 就是专门为 PDF 增加的向下挤压距离！
                        # 如果不够，改大(比如 Pt(40))；如果挤到了第二页，改小(比如 Pt(20))
                        p.paragraph_format.space_before = Pt(10)
                        break

                # 3. 把“加了料”的文档保存为一个专门用来转 PDF 的临时过渡文件
                temp_pdf_docx_path = os.path.join(temp_dir, f"{safe_emp}_for_pdf.docx")
                doc.save(temp_pdf_docx_path)

                # 4. 让 LibreOffice 去转换这个专门为 PDF 定制的过渡文件
                subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', temp_dir, temp_pdf_docx_path
                ], stdout=subprocess.PIPE, stderr=subprocess.PIPE)

                # 5. 将定制生成的 PDF 和最初那个正常的 Word 打包进 ZIP
                temp_pdf_path = os.path.join(temp_dir, f"{safe_emp}_for_pdf.pdf")
                if os.path.exists(temp_pdf_path):
                    with open(temp_pdf_path, 'rb') as f:
                        # 写入压缩包时，把名字变回正常的，假装什么都没发生过
                        zip_file.writestr(f"PDF_Version/{safe_vessel}/{safe_emp}.pdf", f.read())

                with open(temp_docx_path, 'rb') as f:
                    zip_file.writestr(f"Word_Version/{safe_vessel}/{safe_emp}.docx", f.read())

    zip_buffer.seek(0)
    return zip_buffer

# --- 3. Login UI ---
def login_ui():
    _, col_logo, _ = st.columns([2, 1, 2])
    with col_logo:
        try:
            st.image("TSM_Logo.png", use_container_width=True)
        except:
            pass
    st.markdown("<h2 style='text-align: center;'>TSM Ship Information Aggregation System</h2>", unsafe_allow_html=True)
    with st.form("login_form"):
        u_in = st.text_input("User Name")
        p_in = st.text_input("Password", type="password")
        if st.form_submit_button("Log In", use_container_width=True):
            with get_engine().connect() as conn:
                res = conn.execute(text("SELECT role FROM users WHERE username = :u AND password = :p"),
                                   {"u": u_in, "p": p_in}).fetchone()
                if res:
                    st.session_state.clear()
                    st.session_state.logged_in = True
                    st.session_state.username = u_in
                    st.session_state.role = res[0]
                    st.rerun()
                else:
                    st.error("Verification Failed. Please check your credentials.")

if not st.session_state.logged_in:
    login_ui()
    st.stop()


# --- 4. Sidebar ---
st.sidebar.title(f"User: {st.session_state.username}")
if st.sidebar.button("Log Out Safely"):
    st.session_state.clear()
    st.rerun()


# --- 5. Data Retrieval & Tabs ---
@st.cache_data(ttl=60)
def get_ships_list(role, user):
    with get_engine().connect() as conn:
        if role == 'admin':
            return pd.read_sql_query(text("SELECT id, ship_name FROM ships ORDER BY ship_name"), conn)
        return pd.read_sql_query(text("SELECT id, ship_name FROM ships WHERE manager_name = :u ORDER BY ship_name"),
                                 conn, params={"u": user})

ships_df = get_ships_list(st.session_state.role, st.session_state.username)

t_labels = ["Filling & Querying"]
if st.session_state.role == 'admin': t_labels.append("Admin Console")
t_labels.append("Report Center")
tabs = st.tabs(t_labels)

# --- Tab 1: Operations ---
with tabs[0]:
    if ships_df.empty:
        st.warning("No vessels have been assigned yet.")
    else:
        selected_ship = st.selectbox("Select a vessel", ships_df['ship_name'].tolist(), index=st.session_state.ship_index)
        ship_id = int(ships_df[ships_df['ship_name'] == selected_ship]['id'].iloc[0])
        st.divider()

        col_hist, col_input = st.columns([1.2, 1])

        # A. History Record
        with col_hist:
            st.subheader("History Record")

            if st.session_state.confirm_del_id:
                st.warning(f"Prepare to delete the record. (ID: {st.session_state.confirm_del_id})")
                d_col1, d_col2 = st.columns(2)
                with d_col1:
                    if st.button("Confirm deletion", key="confirm_real_del"):
                        with get_engine().begin() as conn:
                            conn.execute(text("DELETE FROM reports WHERE id = :id"),
                                         {"id": st.session_state.confirm_del_id})
                        st.session_state.confirm_del_id = None
                        st.success("The record has been permanently deleted.")
                        time.sleep(1)
                        st.rerun()
                with d_col2:
                    if st.button("Cancel Delete", key="cancel_real_del"):
                        st.session_state.confirm_del_id = None
                        st.rerun()
                st.divider()

            with get_engine().connect() as conn:
                h_df = pd.read_sql_query(text(
                    "SELECT id, report_date, this_week_issue, remarks FROM reports WHERE ship_id = :sid AND is_deleted_by_user = FALSE ORDER BY report_date DESC LIMIT 10"),
                    conn, params={"sid": ship_id})

            if not h_df.empty:
                for idx, row in h_df.iterrows():
                    is_editing = st.session_state.editing_id == row['id']
                    with st.expander(f"{row['report_date']} Content Details", expanded=is_editing):
                        if is_editing:
                            new_val = st.text_area("Modifications:", value=row['this_week_issue'], key=f"ed_{row['id']}")
                            if st.button("Save Updates", key=f"save_{row['id']}"):
                                with get_engine().begin() as conn:
                                    conn.execute(text("UPDATE reports SET this_week_issue = :t WHERE id = :id"),
                                                 {"t": new_val, "id": row['id']})
                                st.session_state.editing_id = None
                                st.rerun()
                        else:
                            raw_content = row['this_week_issue']
                            clean_lines = [re.sub(r'^\d+[\.、\s]*', '', l.strip()) for l in raw_content.split('\n') if l.strip()]
                            st.text("\n".join([f"{i + 1}. {text}" for i, text in enumerate(clean_lines)]))

                            cb1, cb2 = st.columns(2)
                            with cb1:
                                if st.button("Modify", key=f"eb_{row['id']}"):
                                    st.session_state.editing_id = row['id']
                                    st.rerun()
                            with cb2:
                                if st.button("Delete", key=f"db_{row['id']}"):
                                    st.session_state.confirm_del_id = row['id']
                                    st.rerun()
            else:
                st.info("The vessel has no history.")

        # B. Data Input (Indentation fixed to align perfectly with col_hist)
        with col_input:
            st.subheader(f"Fill in - {selected_ship}")

            def handle_submit(sid):
                latest_issue = st.session_state.get(f"ta_{sid}", "")
                latest_remark = st.session_state.get(f"rem_{sid}", "")

                if latest_issue.strip():
                    with get_engine().begin() as conn:
                        conn.execute(text(
                            "INSERT INTO reports (ship_id, report_date, this_week_issue, remarks) VALUES (:sid, :dt, :iss, :rem)"),
                            {"sid": sid, "dt": datetime.now().date(), "iss": latest_issue, "rem": latest_remark})

                    st.session_state[f"ta_{sid}"] = ""
                    st.session_state[f"rem_{sid}"] = ""
                    st.session_state.drafts[sid] = ""
                    st.toast(f"{selected_ship} Data submission successful!")

            if st.button("Import information about the ship from last week.", key=f"import_{ship_id}", use_container_width=True):
                with get_engine().connect() as conn:
                    last_rec = conn.execute(text(
                        "SELECT this_week_issue FROM reports WHERE ship_id = :sid AND is_deleted_by_user = FALSE ORDER BY report_date DESC LIMIT 1"),
                        {"sid": ship_id}).fetchone()
                    if last_rec:
                        st.session_state[f"ta_{ship_id}"] = last_rec[0]
                        st.success("The latest content has been loaded; you can continue editing.")
                        time.sleep(0.5)
                        st.rerun()
                    else:
                        st.warning("No history found.")

            if f"ta_{ship_id}" not in st.session_state:
                st.session_state[f"ta_{ship_id}"] = ""
            if f"rem_{ship_id}" not in st.session_state:
                st.session_state[f"rem_{ship_id}"] = ""

            st.text_area("This week's issue (one line per issue):", height=350, key=f"ta_{ship_id}")
            st.text_input("Remarks (optional)", key=f"rem_{ship_id}")

            st.button(
                "Submit Information",
                use_container_width=True,
                on_click=handle_submit,
                args=(ship_id,)
            )

        st.divider()
        n1, _, n3 = st.columns([1, 4, 1])
        with n1:
            if st.button("Previous"):
                st.session_state.ship_index = (st.session_state.ship_index - 1) % len(ships_df)
                st.rerun()
        with n3:
            if st.button("Next"):
                st.session_state.ship_index = (st.session_state.ship_index + 1) % len(ships_df)
                st.rerun()


# --- Tab 2: Admin Console ---
if st.session_state.role == 'admin':
    with tabs[1]:
        st.subheader("Global Management View")

        m_df = pd.read_sql_query(text("""
            SELECT r.id, s.manager_name as "Manager", s.ship_name as "Vessel", 
                   r.report_date as "Date", r.this_week_issue as "Content"
            FROM reports r JOIN ships s ON r.ship_id = s.id 
            ORDER BY r.report_date DESC
        """), get_engine())

        if not m_df.empty:
            m_df.insert(0, "Select", False)
            ed_df = st.data_editor(m_df, hide_index=True, use_container_width=True)

            to_del = ed_df[ed_df["Select"] == True]["id"].tolist()
            if to_del and st.button("Delete Selected Records"):
                with get_engine().begin() as conn:
                    conn.execute(text("DELETE FROM reports WHERE id IN :ids"), {"ids": tuple(to_del)})
                st.success(f"Successfully deleted {len(to_del)} records.")
                st.rerun()
        else:
            st.info("No global report data available.")

        # =========================================================
        # ✅ 工资单自动化生成区 (内港与外港彻底分离)
        # =========================================================
        st.write("---")
        st.subheader("Automated Paylist Generator")

        # 使用单选按钮让用户明确选择操作模式
        paylist_mode = st.radio(
            "Select Paylist Type:",
            ["In Port Paylist", "Out Port Paylist"],
            horizontal=True
        )

        st.write("")  # 留一行空白让界面更透气

        # ---------------------------------------------------------
        # 模式 A: 内港 (In Port) - 仅生成 Word
        # ---------------------------------------------------------
        if paylist_mode == "In Port Paylist":
            st.info("In Port Mode: Generates Word documents.")

            # 注意：加入了 key 参数，防止上传框冲突
            uploaded_in_port = st.file_uploader("Upload 'SUM-SAL' Excel file (In Port)", type=["xlsx"], key="upload_in")

            if uploaded_in_port is not None:
                if st.button("Generate In Port Paylists (Word ZIP)", use_container_width=True):
                    with st.spinner("Processing In Port documents... Please wait."):
                        try:
                            uploaded_in_port.seek(0)
                            # 调用基础版函数
                            zip_data_in = generate_paylist_zip(uploaded_in_port)
                            st.success("Successfully generated In Port payslips!")

                            st.download_button(
                                label="📥 Download In Port Payslips (.zip)",
                                data=zip_data_in,
                                file_name=f"In_Port_Paylists_{datetime.now().strftime('%Y%m%d')}.zip",
                                mime="application/zip",
                                use_container_width=True
                            )
                        except Exception as e:
                            st.error(f"Error generating In Port paylists: {e}")

        # ---------------------------------------------------------
        # 模式 B: 外港 (Out Port) - 动态计算 + 生成 Word 和 PDF
        # ---------------------------------------------------------
        else:
            st.info(
                "Out Port Mode: generates BOTH Word and PDF documents.")

            uploaded_out_port = st.file_uploader("Upload 'SUM-SAL' Excel file (Out Port)", type=["xlsx"],
                                                 key="upload_out")

            if uploaded_out_port is not None:
                if st.button("Generate Out Port Paylists (Word & PDF ZIP)", use_container_width=True):
                    with st.spinner("Processing calculations and converting PDFs via LibreOffice... Please wait."):
                        try:
                            uploaded_out_port.seek(0)
                            # 调用进阶版函数
                            zip_data_out = generate_advanced_paylist_zip(uploaded_out_port)
                            st.success("Successfully generated Out Port Word & PDF payslips!")

                            st.download_button(
                                label="Download Out Port Payslips (.zip)",
                                data=zip_data_out,
                                file_name=f"Out_Port_Paylists_{datetime.now().strftime('%Y%m%d')}.zip",
                                mime="application/zip",
                                use_container_width=True
                            )
                        except Exception as e:
                            st.error(f"Error generating Out Port paylists: {e}")
# --- Tab 3: Report Center ---
with tabs[-1]:
    st.subheader("Automated Information Preview & Export")

    c1, c2 = st.columns(2)
    with c1:
        start_d = st.date_input("Start Date", value=datetime.now() - timedelta(days=7), key="rep_start")
    with c2:
        end_d = st.date_input("End Date", value=datetime.now(), key="rep_end")

    with get_engine().connect() as conn:
        query = """
                SELECT r.report_date as "Date", s.ship_name as "Vessel", 
                       r.this_week_issue as "Report Content", s.manager_name as "Manager"
                FROM reports r 
                JOIN ships s ON r.ship_id = s.id
                WHERE r.report_date BETWEEN :s AND :e 
                AND r.is_deleted_by_user = FALSE
            """
        params = {"s": start_d, "e": end_d}

        if st.session_state.role != 'admin':
            query += " AND s.manager_name = :u"
            params["u"] = st.session_state.username

        query += " ORDER BY r.report_date DESC"
        export_df = pd.read_sql_query(text(query), conn, params=params)

    st.write("---")
    if st.button("Search and preview records within the selected date", use_container_width=True):
        if not export_df.empty:
            st.success(f"Found {len(export_df)} records.")

            preview_df = export_df.copy()

            def preview_clean(text):
                lines = [re.sub(r'^\d+[\.、\s]*', '', l.strip()) for l in str(text).split('\n') if l.strip()]
                return "\n".join([f"{i + 1}. {t}" for i, t in enumerate(lines)])

            preview_df["Report Content"] = preview_df["Report Content"].apply(preview_clean)

            st.dataframe(
                preview_df,
                use_container_width=True,
                hide_index=True,
                column_config={
                    "Report Content": st.column_config.TextColumn("Detailed Information (Automatically Numbered)", width="large"),
                    "Date": st.column_config.DateColumn("Date")
                }
            )
        else:
            st.warning("No reporting records were found within that date range.")

    st.write("---")

    if not export_df.empty:
        excel_prep_df = export_df.rename(columns={
            "Manager": "manager_name",
            "Vessel": "ship_name",
            "Report Content": "this_week_issue"
        })

        bc1, bc2 = st.columns(2)
        with bc1:
            excel_bin = generate_custom_excel(excel_prep_df)
            st.download_button(
                label="Download Excel Report",
                data=excel_bin,
                file_name=f"Trust_Ship_Report_{start_d}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True
            )

        if st.session_state.role == 'admin':
            with bc2:
                if st.button("Generate PPT Summary Preview", use_container_width=True):
                    ppt_bin = create_ppt_report(excel_prep_df, start_d, end_d)

                    st.download_button(
                        label="Click to Download PPT File",
                        data=ppt_bin,
                        file_name=f"Ship_Meeting_{start_d}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
    else:
        st.info("There is currently no data available for you to view within this date range.")
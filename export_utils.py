import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import sqlalchemy
from sqlalchemy import text
import streamlit as st
from datetime import datetime


# 1. 数据库连接函数（从 Streamlit Secrets 读取）
def get_conn():
    try:
        db_url = st.secrets["postgres_url"]
        # 处理可能出现的协议头不兼容问题
        if db_url.startswith("postgres://"):
            db_url = db_url.replace("postgres://", "postgresql://", 1)
        engine = sqlalchemy.create_engine(db_url)
        return engine.connect()
    except Exception as e:
        st.error(f"导出工具连接数据库失败: {e}")
        return None


# 2. 核心数据抓取：本周问题 + 自动关联上周问题
def get_report_data():
    conn = get_conn()
    if not conn:
        return pd.DataFrame()

    # SQL 逻辑：抓取过去 7 天内提交的所有报告
    main_query = text("""
        SELECT r.id, s.ship_name, s.manager_name, r.report_date, r.this_week_issue, r.remarks, r.ship_id
        FROM reports r
        JOIN ships s ON r.ship_id = s.id
        WHERE r.report_date >= CURRENT_DATE - INTERVAL '7 days'
        ORDER BY r.report_date DESC
    """)

    this_week_records = conn.execute(main_query).fetchall()

    final_data = []
    for row in this_week_records:
        # 子查询：为当前这艘船寻找“比本条记录日期更早”的最新一条记录
        last_week_query = text("""
            SELECT this_week_issue FROM reports 
            WHERE ship_id = :sid AND report_date < :rdate
            ORDER BY report_date DESC LIMIT 1
        """)
        last_res = conn.execute(last_week_query, {"sid": row.ship_id, "rdate": row.report_date}).fetchone()
        last_issue = last_res[0] if last_res else "无历史记录"

        final_data.append({
            "日期": row.report_date,
            "船名": row.ship_name,
            "船舶管理人": row.manager_name,
            "上一周问题": last_issue,
            "本周问题": row.this_week_issue,
            "备注": row.remarks
        })

    conn.close()
    return pd.DataFrame(final_data)


# 3. 生成 Excel
def generate_excel(df, filename):
    df.to_excel(filename, index=False, engine='openpyxl')
    return filename


# 4. 生成 PPT (带颜色标注和自动排版)
def generate_ppt(df, filename):
    prs = Presentation()

    # 如果没有数据，生成一张空白提醒页
    if df.empty:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.text = "本周暂无填报数据"
        prs.save(filename)
        return filename

    for _, row in df.iterrows():
        # 使用“标题和内容”布局
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # 标题：船名
        slide.shapes.title.text = f"🚢 {row['船名']} 会议汇报"

        # 内容正文
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True

        # 第一行：基础信息
        p = tf.paragraphs[0]
        p.text = f"汇报人：{row['船舶管理人']} | 日期：{row['日期']}"
        p.font.size = Pt(18)

        # 第二行：上周回顾
        p = tf.add_paragraph()
        p.text = "\n[上周问题回溯]"
        p.font.bold = True
        p.font.size = Pt(16)

        p = tf.add_paragraph()
        p.text = str(row['上一周问题'])
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(100, 100, 100)  # 灰色表示过去

        # 第三行：本周重点 (醒目红色)
        p = tf.add_paragraph()
        p.text = "\n[本周存在问题]"
        p.font.bold = True
        p.font.size = Pt(18)

        p = tf.add_paragraph()
        p.text = str(row['本周问题'])
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 0, 0)  # 醒目红

        # 第四行：备注
        if row['备注']:
            p = tf.add_paragraph()
            p.text = f"\n备注：{row['备注']}"
            p.font.size = Pt(14)
            p.font.italic = True

    prs.save(filename)
    return filename